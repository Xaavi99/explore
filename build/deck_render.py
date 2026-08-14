"""Assemble a slide plan into PDF + image-PPTX + editable-PPTX.

- PDF and image-PPTX consume the pixel-perfect PNGs from deck_layouts (identical
  to each other by construction).
- editable-PPTX is a NATIVE python-pptx build (real text/shapes/pictures) reading
  the same slide plan + theme; it is close to, not pixel-identical with, the image
  outputs, and its text is editable (fonts depend on the viewer's machine).
"""
import os, tempfile
import pymupdf
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import theme
import deck_layouts as L

LW, LH = 1280, 720
SLIDE_IN_W, SLIDE_IN_H = 13.333, 7.5


def IX(u):
    return Inches(u / LW * SLIDE_IN_W)


def IY(u):
    return Inches(u / LH * SLIDE_IN_H)


def FS(u):
    return Pt(u * 0.75)  # logical (1280-space) size -> PowerPoint points


def C(hex_str):
    return RGBColor(*theme.rgb(hex_str))


# ---------- pixel-perfect PNG rendering ----------
def render_pngs(plan, outdir, quality=88):
    """Render each slide to a JPEG (small, email-friendly; text stays crisp at 2x)."""
    os.makedirs(outdir, exist_ok=True)
    paths = []
    for i, rec in enumerate(plan):
        img = L.render_slide(rec)
        p = os.path.join(outdir, f"slide_{i + 1:02d}.jpg")
        img.save(p, format="JPEG", quality=quality, optimize=True, progressive=True)
        paths.append(p)
    return paths


# ---------- PDF (pymupdf) ----------
def build_pdf(png_paths, out_path):
    doc = pymupdf.open()
    for p in png_paths:
        page = doc.new_page(width=LW, height=LH)
        page.insert_image(pymupdf.Rect(0, 0, LW, LH), filename=p)
    doc.save(out_path, deflate=True)
    doc.close()


# ---------- image-slide PPTX ----------
def build_image_pptx(png_paths, out_path):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_IN_W)
    prs.slide_height = Inches(SLIDE_IN_H)
    blank = prs.slide_layouts[6]
    for p in png_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(p, 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(out_path)


# ---------- native editable PPTX ----------
def _bg(slide, hex_str):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C(hex_str)


def _rect(slide, x, y, w, h, hex_str):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IX(x), IY(y), IX(w), IY(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = C(hex_str)
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _cropped(photo_path, wu, hu, tmpdir, tag):
    """Crop a source photo to the box aspect and return a temp PNG path."""
    if not photo_path or not os.path.exists(photo_path):
        img = L.placeholder(int(wu * 2), int(hu * 2), theme.BRAND_2)
    else:
        img = L.cover_crop(L.load_photo(photo_path), int(wu * 2), int(hu * 2))
    p = os.path.join(tmpdir, f"crop_{tag}.jpg")
    img.save(p, format="JPEG", quality=88, optimize=True, progressive=True)
    return p


def _text(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(IX(x), IY(y), IX(w), IY(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, fname, size, hexc, bold, spacing) in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_before = Pt(0); para.space_after = Pt(spacing)
        run = para.add_run(); run.text = txt
        run.font.name = fname; run.font.size = FS(size)
        run.font.bold = bold; run.font.color.rgb = C(hexc)
    return tb


SERIF_N = "Palatino Linotype"
SANS_N = "Segoe UI"
M = 74  # logical margin (~0.058*1280)
LW, LH = 1280, 720


def _faint_logo(dark, tmpdir):
    """A reduced-opacity copy of the logo for use as a watermark picture."""
    from PIL import Image as PImage
    p = os.path.join(tmpdir, "wm_" + ("d" if dark else "l") + ".png")
    if not os.path.exists(p):
        im = PImage.open(L.LOGO_DARK if dark else L.LOGO_LIGHT).convert("RGBA")
        im.putalpha(im.split()[3].point(lambda v: int(v * 0.5)))
        im.save(p)
    return p


def _native_branding(slide, dark, tmpdir):
    """Footer (company name) + watermark logo on a native editable slide."""
    tb = slide.shapes.add_textbox(IX(0), IY(652), IX(LW), IY(36))
    tf = tb.text_frame; tf.word_wrap = False
    para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
    run = para.add_run(); run.text = L.FOOTER_TEXT
    run.font.name = SANS_N; run.font.size = Pt(7); run.font.bold = True
    run.font.color.rgb = C("FFFFFF") if dark else C(theme.INK_FAINT)
    from PIL import Image as PImage
    lp = _faint_logo(dark, tmpdir)
    im = PImage.open(lp); lh = 0.32; lw = im.width * lh / im.height
    slide.shapes.add_picture(lp, Inches(SLIDE_IN_W - 0.62 - lw), Inches(SLIDE_IN_H - 0.58), Inches(lw), Inches(lh))


def _native_sections(prs, s, rec, tmpdir, idx):
    """Editable region card: collage image column + text column (one flowing box)."""
    side = rec.get("img_side", "left")
    imgw = LW * 0.40
    coll = L._collage(rec.get("left_paths", []), int(imgw * 2), int(LH * 2))
    cp = os.path.join(tmpdir, f"coll{idx}.jpg")
    coll.save(cp, format="JPEG", quality=88, optimize=True, progressive=True)
    if side == "left":
        s.shapes.add_picture(cp, 0, 0, IX(imgw), prs.slide_height)
        tx0, tx1 = imgw + LW * 0.055, LW - M
    else:
        s.shapes.add_picture(cp, IX(LW - imgw), 0, IX(imgw), prs.slide_height)
        tx0, tx1 = M, LW - imgw - LW * 0.05
    maxw = tx1 - tx0
    tcolor = theme.BRAND_2 if rec.get("title_brand") else theme.INK
    paras = [
        (rec["eyebrow"].upper(), SANS_N, 12, theme.GOLD, True, 8),
        (rec["title"], SERIF_N, rec.get("title_pt", 34), tcolor, True, 12),
    ]
    if rec.get("intro"):
        paras.append((rec["intro"], SANS_N, 14, theme.INK_SOFT, False, 14))
    for label, body in rec["sections"]:
        paras.append((label.upper(), SANS_N, 11.5, theme.GOLD, True, 3))
        paras.append((body, SANS_N, 13, theme.INK, True, 12))
    _text(s, tx0, LH * 0.12, maxw, LH * 0.78, paras)


def _native_cta(prs, s, rec, tmpdir, idx):
    panelw = LW * 0.52
    s.shapes.add_picture(_cropped(rec.get("image"), LW - panelw, LH, tmpdir, f"cta{idx}"),
                         IX(panelw), 0, IX(LW - panelw), prs.slide_height)
    _rect(s, 0, 0, panelw, LH, theme.BRAND_2)
    paras = [
        (rec["eyebrow"].upper(), SANS_N, 12, "#e9c983", True, 12),
        (rec["title"], SERIF_N, 40, "#ffffff", False, 14),
    ]
    for text, strong in rec["lines"]:
        paras.append((text, SANS_N, 15 if strong else 14, "#ffffff" if strong else "#cfe0ee", strong, 6))
    _text(s, M, LH * 0.24, panelw - 2 * M, LH * 0.5, paras)


def _slide_native(prs, rec, tmpdir, idx):
    """Keep itinerary-specific slides editable; rasterize fixed brand slides."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    lay = rec["layout"]

    if lay == "cover":
        s.shapes.add_picture(_cropped(rec.get("image"), LW, LH, tmpdir, f"cov{idx}"), 0, 0,
                             width=prs.slide_width, height=prs.slide_height)
        _rect(s, 0, LH * 0.55, LW, LH * 0.45, "#0a1520")
        _text(s, M, LH * 0.60, LW * 0.7, LH * 0.32, [
            (rec["kicker"].upper(), SANS_N, 13, theme.GOLD, True, 8),
            (rec["title"], SERIF_N, 60, "#ffffff", False, 6),
            (rec.get("caption", ""), SANS_N, 16, "#eef2f6", False, 0),
        ])
        _native_branding(s, True, tmpdir)
        return s

    if lay == "sections" and rec.get("editable"):
        _bg(s, theme.PAPER)
        _native_sections(prs, s, rec, tmpdir, idx)
        _native_branding(s, False, tmpdir)
        return s

    if lay == "cta":
        _native_cta(prs, s, rec, tmpdir, idx)
        _native_branding(s, True, tmpdir)
        return s

    # fixed brand slides: reuse the pixel-perfect Pillow render (branding baked in)
    png = os.path.join(tmpdir, f"img_{idx}.jpg")
    L.render_slide(rec).save(png, format="JPEG", quality=88, optimize=True, progressive=True)
    s.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)
    return s


def build_editable_pptx(plan, out_path, tmpdir):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_IN_W)
    prs.slide_height = Inches(SLIDE_IN_H)
    for i, rec in enumerate(plan):
        _slide_native(prs, rec, tmpdir, i)
    prs.save(out_path)


# ---------- orchestration ----------
def generate(plan, slug, title, distroot):
    outdir = os.path.join(distroot, slug)
    png_dir = os.path.join(outdir, "slides")
    os.makedirs(png_dir, exist_ok=True)
    pngs = render_pngs(plan, png_dir)

    pdf = os.path.join(outdir, f"{slug}.pdf")
    pptx_img = os.path.join(outdir, f"{slug}.pptx")
    pptx_edit = os.path.join(outdir, f"{slug}-editable.pptx")
    build_pdf(pngs, pdf)
    build_image_pptx(pngs, pptx_img)
    with tempfile.TemporaryDirectory() as td:
        build_editable_pptx(plan, pptx_edit, td)
    return {"pdf": pdf, "pptx": pptx_img, "pptx_editable": pptx_edit, "slides": pngs}
